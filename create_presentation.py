#!/usr/bin/env python3
"""Create IMGC PAS Client Demo Presentation using python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SCREENSHOTS_DIR = os.path.join(os.path.dirname(__file__), 'screenshots')
OUTPUT_FILE = os.path.join(os.path.dirname(__file__), 'IMGC_PAS_Client_Demo.pptx')

# Brand Colors
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
DARK_NAVY = RGBColor(0x0F, 0x1B, 0x2D)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
MEDIUM_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
ORANGE_LIGHT = RGBColor(0xFF, 0xF7, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color=WHITE):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """Add a rounded rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4)):
    """Add paragraph to existing text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_points(slide, left, top, width, height, items, font_size=13, color=DARK_TEXT):
    """Add bullet point list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        p.level = 0
    return txBox


def add_orange_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6)):
    """Add orange accent bar"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    return shape


def add_screenshot(slide, filename, left, top, width, height=None):
    """Add screenshot image"""
    filepath = os.path.join(SCREENSHOTS_DIR, filename)
    if os.path.exists(filepath):
        if height:
            pic = slide.shapes.add_picture(filepath, left, top, width, height)
        else:
            pic = slide.shapes.add_picture(filepath, left, top, width)
        return pic
    return None


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, WHITE)

# Navy header bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.2), fill_color=NAVY)

# Orange accent line
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), fill_color=ORANGE)

# Title text
add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.6),
             'IMGC CORPORATION', font_size=16, color=ORANGE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
             'Policy Administration System', font_size=42, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.5),
             'Enterprise Platform — Complete Application Walkthrough', font_size=20, color=RGBColor(0xCB, 0xD5, 0xE1))
add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.4),
             'Version 3.2.1  |  Confidential  |  March 2026', font_size=14, color=MEDIUM_GRAY)

# Feature highlights below
features = [
    ('14+', 'Modules'),
    ('150+', 'API Endpoints'),
    ('5', 'User Roles'),
    ('256-bit', 'Encryption'),
]
for i, (num, label) in enumerate(features):
    x = Inches(1.5 + i * 2.8)
    add_rounded_rect(slide, x, Inches(3.8), Inches(2.2), Inches(1.2), fill_color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.2), Inches(3.95), Inches(1.8), Inches(0.5),
                 num, font_size=32, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.45), Inches(1.8), Inches(0.4),
                 label, font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Platform capabilities
capabilities = [
    'Policy Lifecycle Management — Quote to issuance in minutes',
    'Claims Automation — FNOL to settlement tracking',
    'Risk Intelligence — Real-time underwriting scoring',
    'Billing & Payments — Full invoicing and payment lifecycle',
    'Compliance Ready — Regulatory audit trails built-in (SOC 2, IRDAI)',
]
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             'Platform Capabilities', font_size=18, color=NAVY, bold=True)
add_bullet_points(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(1.5),
                  ['  ' + c for c in capabilities], font_size=13, color=DARK_TEXT)


# ============================================================================
# SLIDE 2: Login & Authentication
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill_color=NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(4), Inches(0.5),
             '01  LOGIN & AUTHENTICATION', font_size=22, color=ORANGE, bold=True)

# Screenshot
add_screenshot(slide, '01-login.png', Inches(0.4), Inches(1.1), Inches(7.5))

# Right panel - features
rx = Inches(8.2)
add_orange_accent_bar(slide, rx, Inches(1.2), height=Inches(0.5))
add_text_box(slide, rx + Inches(0.2), Inches(1.2), Inches(4.5), Inches(0.5),
             'Secure Authentication', font_size=18, color=NAVY, bold=True)

auth_features = [
    '• JWT-based authentication with refresh tokens',
    '• Encrypted password storage (bcrypt)',
    '• Session management with auto-expiry',
    '• Role-based access control from login',
    '• 256-bit encryption end-to-end',
]
add_bullet_points(slide, rx + Inches(0.2), Inches(1.8), Inches(4.8), Inches(2),
                  auth_features, font_size=12, color=DARK_TEXT)

add_orange_accent_bar(slide, rx, Inches(3.6), height=Inches(0.5))
add_text_box(slide, rx + Inches(0.2), Inches(3.6), Inches(4.5), Inches(0.5),
             'Enterprise UX Design', font_size=18, color=NAVY, bold=True)

ux_features = [
    '• Split-screen: Brand showcase + Login form',
    '• Platform highlights (Policy, Claims, Risk, Compliance)',
    '• SOC 2 & IRDAI compliance badges',
    '• Responsive design for all devices',
    '• Orange accent theme (30% coverage)',
]
add_bullet_points(slide, rx + Inches(0.2), Inches(4.2), Inches(4.8), Inches(2),
                  ux_features, font_size=12, color=DARK_TEXT)

add_orange_accent_bar(slide, rx, Inches(5.8), height=Inches(0.5))
add_text_box(slide, rx + Inches(0.2), Inches(5.8), Inches(4.5), Inches(0.5),
             '5 User Roles', font_size=18, color=NAVY, bold=True)

roles = [
    '• Admin — Full system access',
    '• Underwriter — Risk assessment & approval',
    '• Claims Adjuster — Claims processing',
    '• Agent — Customer-facing operations',
    '• Auditor — Read-only compliance review',
]
add_bullet_points(slide, rx + Inches(0.2), Inches(6.4), Inches(4.8), Inches(1),
                  roles, font_size=12, color=DARK_TEXT)


# ============================================================================
# Helper function for module slides
# ============================================================================
def create_module_slide(slide_num, title, screenshot_file, features_left, features_right, left_title="Key Features", right_title="Capabilities"):
    """Create a standard module slide with screenshot + feature descriptions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill_color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 f'{slide_num:02d}  {title.upper()}', font_size=22, color=ORANGE, bold=True)

    # Screenshot
    add_screenshot(slide, screenshot_file, Inches(0.3), Inches(1.1), Inches(7.8))

    # Right panel
    rx = Inches(8.4)

    # Left feature box
    add_orange_accent_bar(slide, rx, Inches(1.2), height=Inches(0.5))
    add_text_box(slide, rx + Inches(0.2), Inches(1.2), Inches(4.5), Inches(0.5),
                 left_title, font_size=18, color=NAVY, bold=True)
    add_bullet_points(slide, rx + Inches(0.2), Inches(1.8), Inches(4.6), Inches(2.5),
                      features_left, font_size=12, color=DARK_TEXT)

    # Right feature box
    y_offset = 1.8 + len(features_left) * 0.28 + 0.3
    if y_offset < 3.8:
        y_offset = 3.8
    add_orange_accent_bar(slide, rx, Inches(y_offset), height=Inches(0.5))
    add_text_box(slide, rx + Inches(0.2), Inches(y_offset), Inches(4.5), Inches(0.5),
                 right_title, font_size=18, color=NAVY, bold=True)
    add_bullet_points(slide, rx + Inches(0.2), Inches(y_offset + 0.6), Inches(4.6), Inches(2.5),
                      features_right, font_size=12, color=DARK_TEXT)

    return slide


# ============================================================================
# SLIDE 3: Executive Dashboard (Top)
# ============================================================================
create_module_slide(
    2, 'Executive Dashboard',
    '02-dashboard-top.png',
    [
        '• 7 Executive KPIs with target tracking:',
        '  — Loss Ratio (Target: 65%)',
        '  — Expense Ratio (Target: 35%)',
        '  — Combined Ratio (Target: 100%)',
        '  — Retention Rate (Target: 85%)',
        '  — Collection Rate (Target: 90%)',
        '  — Avg Cycle Time (Target: 30 days)',
        '  — Avg Severity (Target: ₹50,000)',
    ],
    [
        '• 4 KPI Summary Cards:',
        '  — Total Policies: 52 (23 active)',
        '  — Total Premium: ₹43,87,475',
        '  — Total Coverage: ₹13,73,00,324',
        '  — Total Claims: 23 (₹1.00Cr total)',
        '• Policies by Status bar chart',
        '• Policy Type Distribution pie chart',
    ],
    left_title='Executive KPIs',
    right_title='KPI Cards & Charts'
)

# ============================================================================
# SLIDE 4: Dashboard Analytics (Bottom)
# ============================================================================
create_module_slide(
    2, 'Dashboard — Analytics & Widgets',
    '03-dashboard-bottom.png',
    [
        '• Risk Distribution pie chart:',
        '  — Low: 24, Medium: 16, High: 10',
        '• Claims by Status horizontal bar chart',
        '• Underwriting Decisions donut chart:',
        '  — Auto-Approve: 28, Refer: 12, Reject: 10',
    ],
    [
        '• My Tasks widget with SLA deadlines',
        '• Underwriting Overview:',
        '  — 50 Total Evaluations, 41.5 Avg Risk Score',
        '• Billing Overview:',
        '  — 10 Accounts, ₹13,058 Outstanding',
        '  — ₹10,033 Collected',
        '• Overdue Accounts alert (5 accounts, ₹2,950)',
    ],
    left_title='Analytics Charts',
    right_title='Dashboard Widgets'
)

# ============================================================================
# SLIDE 5: Policy Management
# ============================================================================
create_module_slide(
    3, 'Policy Management',
    '04-policies.png',
    [
        '• Comprehensive policy table with 8 columns',
        '• 12 Policy statuses: Draft, Quoted, Bound,',
        '  Issued, Active, Endorsed, Renewal Pending,',
        '  Reinstated, Lapsed, Cancelled, Expired',
        '• Filter by status and policy type',
        '• Real-time search across all fields',
        '• Sortable columns (ID, Type, Premium, etc.)',
        '• Pagination: 10 records/page',
    ],
    [
        '• "New Quote" — Start quote wizard',
        '• "New Policy" — Direct policy creation',
        '• 3 Policy Types:',
        '  — Mortgage Guarantee',
        '  — Credit Protection',
        '  — Coverage Plus',
        '• Color-coded risk badges (Low/Medium/High)',
        '• Status badges with icons',
    ],
    left_title='Policy List Management',
    right_title='Quick Actions'
)

# ============================================================================
# SLIDE 6: Policy Detail & Lifecycle
# ============================================================================
create_module_slide(
    4, 'Policy Detail & Lifecycle',
    '05-policy-detail.png',
    [
        '• Visual Lifecycle Stepper:',
        '  Draft → Quoted → Bound → Issued → Active',
        '• Checkmarks on completed stages',
        '• Available Actions display:',
        '  Endorsed, Renewal Pending, Lapsed,',
        '  Cancelled, Expired',
        '• Header buttons: Initiate Renewal,',
        '  Lapse, Cancel',
    ],
    [
        '• 8-Tab Detail View:',
        '  1. Details — Policy & Financial info',
        '  2. Endorsements — Mid-term changes',
        '  3. Renewals — Renewal history',
        '  4. Timeline — Event chronology',
        '  5. Versions — Version comparison',
        '  6. Billing — Payment tracking',
        '  7. Documents — Generated docs',
        '  8. Audit — Complete audit trail',
    ],
    left_title='Policy Lifecycle',
    right_title='Detail Tabs'
)

# ============================================================================
# SLIDE 7: Customer Management
# ============================================================================
create_module_slide(
    5, 'Customer Management',
    '06-customers.png',
    [
        '• Customer table: Name, Email, Phone,',
        '  Status, Policies, Risk Profile',
        '• Real-time search across all fields',
        '• Active/Inactive status badges',
        '• Risk profile indicators (Low/Med/High)',
        '• Direct link to customer policies',
        '• Pagination with record counts',
    ],
    [
        '• Create new customers with full profile',
        '• Edit existing customer information',
        '• View linked policies and claims',
        '• Risk profile assessment per customer',
        '• Customer detail with all linked entities',
        '• Search by name, email, or phone',
    ],
    left_title='Customer Records',
    right_title='Customer Operations'
)

# ============================================================================
# SLIDE 8: Underwriting Engine
# ============================================================================
create_module_slide(
    6, 'Underwriting Engine',
    '07-underwriting.png',
    [
        '• Automated risk scoring (1-100 scale)',
        '• 6 Underwriting Rules:',
        '  — Property value limits',
        '  — LTV ratio thresholds',
        '  — Credit score requirements',
        '  — Employment verification',
        '  — Property type eligibility',
        '  — Coverage ratio limits',
    ],
    [
        '• Decision Categories:',
        '  — Auto-Approve (Low Risk)',
        '  — Refer (Medium Risk)',
        '  — Reject (High Risk)',
        '• Manual override with justification',
        '• Configurable rules (Admin)',
        '• Referral queue for senior UW review',
        '• Delegated authority management',
    ],
    left_title='Risk Assessment',
    right_title='Decisions & Workflow'
)

# ============================================================================
# SLIDE 9: Claims Management
# ============================================================================
create_module_slide(
    7, 'Claims Management',
    '08-claims.png',
    [
        '• Claims table: ID, Policy, Type, Status,',
        '  Amount, Filed Date',
        '• Status flow:',
        '  Filed → Under Review → Approved/Rejected',
        '  → Settled',
        '• FNOL intake wizard (5-step process)',
        '• Reserve management & tracking',
        '• Fraud indicator scoring',
    ],
    [
        '• FNOL Intake Steps:',
        '  1. Incident Information',
        '  2. Policy Lookup',
        '  3. Damage Details',
        '  4. Document Upload',
        '  5. Review & Submit',
        '• Claims timeline with event history',
        '• Loss mitigation tracking',
        '• Settlement processing',
    ],
    left_title='Claims Processing',
    right_title='FNOL Workflow'
)

# ============================================================================
# SLIDE 10: Billing & Payments
# ============================================================================
create_module_slide(
    8, 'Billing & Payments',
    '09-billing.png',
    [
        '• Billing accounts linked to policies',
        '• Auto-generated invoices with line items',
        '• Payment tracking by method:',
        '  Cash, Check, UPI, Bank Transfer',
        '• Installment plans: Annual, Semi-Annual,',
        '  Quarterly, Monthly',
        '• Receivables Aging chart',
        '  (Current, 30d, 60d, 90d, 90+)',
    ],
    [
        '• Grace period automation',
        '  (Overdue → Grace → Lapse)',
        '• Double-entry ledger:',
        '  Premiums, Payments, Refunds, Write-offs',
        '• KPIs:',
        '  — 10 Total Accounts (7 Active)',
        '  — ₹13,058 Outstanding',
        '  — ₹10,033 Collected',
        '  — 5 Overdue (₹2,950)',
    ],
    left_title='Billing Lifecycle',
    right_title='Automation & Ledger'
)

# ============================================================================
# SLIDE 11: Task Management
# ============================================================================
create_module_slide(
    9, 'Task Management',
    '10-tasks.png',
    [
        '• Summary cards: Open, In Progress,',
        '  Completed, Overdue counts',
        '• Task table: Title, Assignee, Priority,',
        '  Due Date, Status',
        '• Priority: Urgent, High, Medium, Low',
        '• SLA countdown with overdue alerts',
        '• Task assignment/reassignment',
    ],
    [
        '• Auto-task creation from workflows:',
        '  — New claim → Claims adjuster task',
        '  — UW referral → Senior UW task',
        '  — Renewal due → Operations task',
        '  — Overdue payment → Billing task',
        '• Dashboard widget for quick access',
        '• Task detail drawer with activity log',
    ],
    left_title='Task Queue',
    right_title='Workflow Automation'
)

# ============================================================================
# SLIDE 12: Document Management
# ============================================================================
create_module_slide(
    10, 'Document Management',
    '11-documents.png',
    [
        '• Category tabs: All, Policy, Claims,',
        '  Underwriting, Correspondence',
        '• Document table: Name, Category,',
        '  Entity, Date, Size, Actions',
        '• Drag-and-drop file upload',
        '• Document generation from templates',
        '• Full-text search',
        '• Version control with history',
    ],
    [
        '• Generated Document Types:',
        '  — Policy Declarations',
        '  — Certificates of Insurance',
        '  — Endorsement Schedules',
        '  — Claims Correspondence',
        '  — Billing Invoices',
        '• In-app preview for PDF/images',
        '• Bulk download as ZIP',
    ],
    left_title='Document System',
    right_title='Generation & Templates'
)

# ============================================================================
# SLIDE 13: Reports & Analytics
# ============================================================================
create_module_slide(
    11, 'Reports & Analytics',
    '12-reports.png',
    [
        '• Report Types:',
        '  — Policy Reports (premium trends)',
        '  — Claims Reports (severity, cycle time)',
        '  — Financial Reports (revenue, loss ratios)',
        '• Interactive bar charts and pie charts',
        '• Summary metrics with trend indicators',
    ],
    [
        '• Date range presets:',
        '  This Month, Last Month, Last Quarter,',
        '  Year-to-Date, Custom range',
        '• Tab-based navigation',
        '  (Policy, Claims, Financial)',
        '• Export capabilities for data analysis',
        '• Visual summaries with KPI cards',
    ],
    left_title='Report Types',
    right_title='Filtering & Export'
)

# ============================================================================
# SLIDE 14: Administration
# ============================================================================
create_module_slide(
    12, 'Administration & User Management',
    '13-admin.png',
    [
        '• User table: Name, Email, Role,',
        '  Status, Last Active, Actions',
        '• RBAC with 5 predefined roles',
        '• Add/Edit/Activate/Deactivate users',
        '• Color-coded role badges',
        '• Complete audit logging',
    ],
    [
        '• Role Permissions Matrix:',
        '  — Admin: Full access to all modules',
        '  — Underwriter: Evaluation access',
        '  — Claims Adjuster: Claims processing',
        '  — Agent: Customer-facing operations',
        '  — Auditor: Read-only compliance',
        '• Session management controls',
    ],
    left_title='User Management',
    right_title='RBAC & Permissions'
)

# ============================================================================
# SLIDE 15: Compliance
# ============================================================================
create_module_slide(
    13, 'Compliance Management',
    '14-compliance.png',
    [
        '• Regulatory requirements tracker',
        '• Status: Compliant, Non-Compliant,',
        '  Under Review, Due',
        '• IRDAI regulatory filing status',
        '• Compliance calendar with deadlines',
        '• Alert system for upcoming deadlines',
    ],
    [
        '• Complete audit trail for all actions',
        '• User action logging with timestamps',
        '• Data change history (before/after)',
        '• Compliance report generation',
        '• SOC 2 documentation',
        '• Regulatory filing tracking',
    ],
    left_title='Regulatory Compliance',
    right_title='Audit & Governance'
)

# ============================================================================
# SLIDE 16: Module Summary
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill_color=NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
             '14  MODULE SUMMARY', font_size=22, color=ORANGE, bold=True)

# Module table
modules = [
    ('1', 'Dashboard', 'Executive KPIs, 7 charts, role-based views, widgets'),
    ('2', 'Policies', 'Full lifecycle (Quote→Bind→Issue→Endorse→Renew→Cancel), versioning'),
    ('3', 'Customers', 'CRUD operations, risk profiles, linked policies/claims'),
    ('4', 'Underwriting', '6 rules, risk scoring (1-100), manual override, referral queue'),
    ('5', 'Claims', 'FNOL intake, status workflow, reserve management, fraud scoring'),
    ('6', 'Billing', 'Invoicing, payments, installments, aging reports, ledger'),
    ('7', 'Tasks', 'Queue management, SLA tracking, auto-assignment from workflows'),
    ('8', 'Documents', 'Upload, template-based generation, categories, versioning'),
    ('9', 'Reports', 'Policy/Claims/Financial reports, date filtering, charts'),
    ('10', 'Admin', 'User CRUD, RBAC (5 roles), audit logs, session management'),
    ('11', 'Compliance', 'Regulatory tracking, deadline alerts, filing status'),
    ('12', 'Renewals', 'Renewal queue, batch processing, auto-renewal workflow'),
    ('13', 'FNOL', 'Structured 5-step intake wizard, document collection'),
    ('14', 'Notifications', 'In-app alerts, bell icon, read/unread management'),
]

# Table header
add_rounded_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.45), fill_color=NAVY)
add_text_box(slide, Inches(0.7), Inches(1.15), Inches(0.6), Inches(0.35), '#', font_size=12, color=ORANGE, bold=True)
add_text_box(slide, Inches(1.4), Inches(1.15), Inches(2), Inches(0.35), 'Module', font_size=12, color=ORANGE, bold=True)
add_text_box(slide, Inches(3.5), Inches(1.15), Inches(8), Inches(0.35), 'Key Capabilities', font_size=12, color=ORANGE, bold=True)

for i, (num, name, desc) in enumerate(modules):
    y = Inches(1.6 + i * 0.38)
    bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
    add_rounded_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.36), fill_color=bg)
    add_text_box(slide, Inches(0.7), y + Inches(0.02), Inches(0.6), Inches(0.3), num, font_size=11, color=ORANGE, bold=True)
    add_text_box(slide, Inches(1.4), y + Inches(0.02), Inches(2), Inches(0.3), name, font_size=11, color=NAVY, bold=True)
    add_text_box(slide, Inches(3.5), y + Inches(0.02), Inches(9), Inches(0.3), desc, font_size=11, color=DARK_TEXT)


# ============================================================================
# SLIDE 17: Technology Stack
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill_color=NAVY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
             '15  TECHNOLOGY STACK & ARCHITECTURE', font_size=22, color=ORANGE, bold=True)

# Tech stack cards
stack = [
    ('Frontend', 'Next.js 14, React 18, TypeScript', 'Server-side rendering, type safety, modern component architecture'),
    ('Styling', 'Tailwind CSS, Lucide Icons', 'Enterprise UI with consistent design system, 30% orange accent'),
    ('Charts', 'Recharts', 'Interactive data visualizations — bar, pie, donut, horizontal bar charts'),
    ('Backend', 'Express.js, TypeScript', 'RESTful API server with 150+ endpoints, middleware architecture'),
    ('Auth', 'JWT + Refresh Tokens', 'Stateless authentication, RBAC, session management'),
    ('Storage', 'JSON File Persistence', 'Rapid prototyping with database-ready architecture'),
]

for i, (title, tech, desc) in enumerate(stack):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.2 + row * 2.0)

    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(1.7), fill_color=LIGHT_GRAY)
    add_orange_accent_bar(slide, x + Inches(0.15), y + Inches(0.2), height=Inches(0.4))
    add_text_box(slide, x + Inches(0.4), y + Inches(0.2), Inches(3.2), Inches(0.4), title, font_size=16, color=NAVY, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.65), Inches(3.2), Inches(0.35), tech, font_size=13, color=ORANGE, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(1.0), Inches(3.2), Inches(0.6), desc, font_size=11, color=MEDIUM_GRAY)

# Non-functional features
nf_title_y = Inches(5.4)
add_text_box(slide, Inches(0.5), nf_title_y, Inches(12), Inches(0.5),
             'Non-Functional Features', font_size=18, color=NAVY, bold=True)

nf_features = [
    ('Responsive', 'Desktop, tablet, mobile'),
    ('Accessible', 'WCAG 2.1, ARIA labels, keyboard nav'),
    ('Performant', 'Skeleton loading, staggered animations'),
    ('Secure', '256-bit encryption, RBAC, JWT, input sanitization'),
    ('Enterprise UX', 'Toasts, confirmations, breadcrumbs, search'),
]

for i, (title, desc) in enumerate(nf_features):
    x = Inches(0.5 + i * 2.5)
    add_rounded_rect(slide, x, Inches(5.9), Inches(2.2), Inches(1.0), fill_color=ORANGE_LIGHT)
    add_text_box(slide, x + Inches(0.15), Inches(5.95), Inches(1.9), Inches(0.35), title, font_size=12, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(6.3), Inches(1.9), Inches(0.55), desc, font_size=10, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 18: Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

# Orange accent line at top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ORANGE)

# Center content
add_text_box(slide, Inches(2), Inches(2.0), Inches(9.3), Inches(0.5),
             'IMGC CORPORATION', font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(2.8), Inches(9.3), Inches(1),
             'Thank You', font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(0.5),
             'Policy Administration System — Enterprise Edition', font_size=20, color=RGBColor(0xCB, 0xD5, 0xE1), alignment=PP_ALIGN.CENTER)

# Divider line
add_shape(slide, Inches(5), Inches(4.6), Inches(3.3), Inches(0.03), fill_color=ORANGE)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
             'For questions or demo access, contact the PAS Development Team', font_size=16, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
             'Confidential  |  IMGC PAS v3.2.1  |  March 2026', font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Feature badges at bottom
badges = ['14+ Modules', '150+ APIs', '5 User Roles', '256-bit Encrypted', 'SOC 2 Compliant', 'IRDAI Regulated']
for i, badge in enumerate(badges):
    x = Inches(1.0 + i * 1.95)
    add_rounded_rect(slide, x, Inches(6.4), Inches(1.7), Inches(0.5), fill_color=RGBColor(0x1E, 0x40, 0x6E))
    add_text_box(slide, x + Inches(0.1), Inches(6.45), Inches(1.5), Inches(0.4),
                 badge, font_size=10, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# Save presentation
prs.save(OUTPUT_FILE)
print(f'Presentation saved to: {OUTPUT_FILE}')
print(f'Total slides: {len(prs.slides)}')
